#!/usr/bin/env python3
"""
上本町ハイハイタウン エリア分析レポート パワーポイント生成スクリプト
不動産コンサルタント 美玖
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================
# 色定義
# ============================================================
PRIMARY_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x1E)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
RED_PRICE = RGBColor(0xE8, 0x3F, 0x3F)
BLUE_PRICE = RGBColor(0x2E, 0x86, 0xDE)
LIGHT_TEAL = RGBColor(0xE0, 0xF2, 0xF1)
LIGHT_GOLD = RGBColor(0xFD, 0xF0, 0xD0)
LIGHT_NAVY = RGBColor(0xE8, 0xEC, 0xF1)
MEDIUM_GRAY = RGBColor(0x90, 0x90, 0x90)
TABLE_HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)
TABLE_ROW_EVEN = RGBColor(0xF5, 0xF5, 0xF0)
TABLE_ROW_ODD = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = "Meiryo"
FONT_NAME_BOLD = "Meiryo"

# ============================================================
# データ読み込み
# ============================================================
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
JSON_PATH = os.path.join(SCRIPT_DIR, "suumo_data_20260325.json")

with open(JSON_PATH, "r", encoding="utf-8") as f:
    suumo_data = json.load(f)

listings = suumo_data["haihaitown_listings"]
summary = suumo_data["summary"]


# ============================================================
# ユーティリティ関数
# ============================================================
def set_slide_bg(slide, color):
    """スライド背景色を設定"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text, slide_num=None):
    """スライド上部にネイビーのヘッダーバーを追加"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_NAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)

    # スライド番号
    if slide_num is not None:
        num_box = slide.shapes.add_textbox(
            Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4)
        )
        tf2 = num_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = str(slide_num)
        p2.font.size = Pt(11)
        p2.font.color.rgb = MEDIUM_GRAY
        p2.font.name = FONT_NAME
        p2.alignment = PP_ALIGN.RIGHT


def add_accent_line(slide, left, top, width):
    """アクセントラインを追加"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_TEAL
    line.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """テキストボックスを追加"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.alignment = alignment
    return box


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=12, text_color=DARK_TEXT, bold=False,
                     alignment=PP_ALIGN.CENTER):
    """角丸四角形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(0.5)

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = FONT_NAME
        p.font.bold = bold
        p.alignment = alignment
    return shape


def add_circle(slide, left, top, size, fill_color, text="",
               font_size=11, text_color=WHITE, bold=True):
    """円を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = FONT_NAME
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow_right(slide, left, top, width, height, color=ACCENT_TEAL):
    """右矢印を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_multiline_text(slide, left, top, width, height, lines,
                       font_size=13, color=DARK_TEXT, line_spacing=1.5,
                       bold=False, alignment=PP_ALIGN.LEFT):
    """複数行テキストボックス"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return box


# ============================================================
# プレゼンテーション生成
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

slide_num_counter = 0


def new_slide():
    global slide_num_counter
    slide_num_counter += 1
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s, LIGHT_BG)
    return s


# ============================================================
# スライド1: 表紙
# ============================================================
slide = new_slide()
set_slide_bg(slide, PRIMARY_NAVY)

# 上部装飾ライン
line_top = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06)
)
line_top.fill.solid()
line_top.fill.fore_color.rgb = ACCENT_GOLD
line_top.line.fill.background()

# メインタイトル
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "🏠 上本町ハイハイタウン", font_size=44, color=WHITE,
             bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "エリア分析レポート", font_size=36, color=ACCENT_GOLD,
             bold=True, alignment=PP_ALIGN.CENTER)

# 区切り線
add_accent_line(slide, Inches(4.5), Inches(3.9), Inches(4.3))

add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.6),
             "不動産コンサルタント  美玖", font_size=22, color=WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
             "2026年3月", font_size=18, color=RGBColor(0xBB, 0xBB, 0xCC),
             alignment=PP_ALIGN.CENTER)

# 下部装飾
line_bot = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06)
)
line_bot.fill.solid()
line_bot.fill.fore_color.rgb = ACCENT_GOLD
line_bot.line.fill.background()


# ============================================================
# スライド2: 目次
# ============================================================
slide = new_slide()
add_header_bar(slide, "📋 目次", slide_num_counter)

toc_items = [
    ("01", "ハイハイタウンとは"),
    ("02", "アクセスの良さ"),
    ("03", "建物スペック"),
    ("04", "周辺施設マップ"),
    ("05", "現在の売出し状況"),
    ("06", "価格帯分布"),
    ("07", "㎡単価の分析"),
    ("08", "築年数と価格の関係"),
    ("09", "賃貸利回り"),
    ("10", "エリアの価格トレンド"),
    ("11", "メリット"),
    ("12", "デメリット・注意点"),
    ("13", "2030年に向けた戦略"),
    ("14", "購入時のチェックポイント"),
    ("15", "出典・参考情報"),
    ("16", "まとめ"),
]

col1_x = Inches(1.0)
col2_x = Inches(7.0)
start_y = Inches(1.4)
row_h = Inches(0.37)

for i, (num, title) in enumerate(toc_items):
    col = 0 if i < 8 else 1
    x = col1_x if col == 0 else col2_x
    y = start_y + row_h * (i if col == 0 else i - 8)

    add_circle(slide, x, y, Inches(0.32), ACCENT_TEAL, num, font_size=9)
    add_text_box(slide, x + Inches(0.45), y + Inches(0.02), Inches(5), Inches(0.32),
                 title, font_size=14, color=DARK_TEXT)


# ============================================================
# スライド3: ハイハイタウンとは
# ============================================================
slide = new_slide()
add_header_bar(slide, "🏢 ハイハイタウンとは", slide_num_counter)

# 左側: 建物概要カード
card = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(5.8),
                        WHITE)
info_lines = [
    "【正式名称】 上本町ハイハイタウン",
    "",
    "【竣工】 1980年9月",
    "",
    "【構造】 鉄骨鉄筋コンクリート造（SRC造）",
    "",
    "【規模】 地上15階・地下2階",
    "",
    "【商業】 地下2階〜4階 = 商業施設",
    "         （100店舗以上の飲食店・商店）",
    "",
    "【住居】 上層階 = 住居（約287戸）",
    "",
    "【所在地】 大阪府大阪市天王寺区",
    "           上本町6丁目3-31",
]
add_multiline_text(slide, Inches(0.9), Inches(1.5), Inches(5.2), Inches(5.5),
                   info_lines, font_size=13, line_spacing=1.2)

# 右側: 建物イメージ図
# フロア構成図
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5), Inches(0.5),
             "🏗️ フロア構成イメージ", font_size=16, bold=True, color=PRIMARY_NAVY)

floors = [
    ("15階", "住居エリア", LIGHT_TEAL, DARK_TEXT),
    ("〜", "（約287戸）", LIGHT_TEAL, DARK_TEXT),
    ("5階", "1LDK〜3LDK", LIGHT_TEAL, DARK_TEXT),
    ("4階", "商業施設", LIGHT_GOLD, DARK_TEXT),
    ("〜", "飲食店", LIGHT_GOLD, DARK_TEXT),
    ("1階", "商店街", LIGHT_GOLD, DARK_TEXT),
    ("地下1階", "商業施設", RGBColor(0xE8, 0xE0, 0xD0), DARK_TEXT),
    ("地下2階", "駐車場等", RGBColor(0xE0, 0xD8, 0xC8), DARK_TEXT),
]

for i, (floor_label, desc, bg, tc) in enumerate(floors):
    y = Inches(1.9) + Inches(0.55) * i
    rect = add_rounded_rect(slide, Inches(7.2), y, Inches(2.0), Inches(0.45),
                            bg, floor_label, font_size=11, text_color=tc, bold=True)
    add_text_box(slide, Inches(9.4), y + Inches(0.05), Inches(3), Inches(0.4),
                 desc, font_size=11, color=tc)


# ============================================================
# スライド4: アクセスの良さ
# ============================================================
slide = new_slide()
add_header_bar(slide, "🚃 アクセスの良さ", slide_num_counter)

# 中央: ハイハイタウン
center_x = Inches(5.4)
center_y = Inches(3.2)
add_circle(slide, center_x, center_y, Inches(1.6), PRIMARY_NAVY,
           "ハイハイ\nタウン", font_size=14, text_color=WHITE)

# 周辺駅・施設を配置
stations = [
    ("近鉄\n大阪上本町駅\n（直結・徒歩1分）", Inches(1.0), Inches(2.0), RED_PRICE),
    ("大阪メトロ\n谷町九丁目駅\n（徒歩3分）", Inches(1.0), Inches(4.8), BLUE_PRICE),
    ("近鉄百貨店\n（隣接）", Inches(9.5), Inches(2.0), ACCENT_GOLD),
    ("上本町\nYUFURA\n（隣接）", Inches(9.5), Inches(4.8), ACCENT_TEAL),
]

for label, x, y, color in stations:
    add_rounded_rect(slide, x, y, Inches(2.5), Inches(1.4), color,
                     label, font_size=12, text_color=WHITE, bold=True)

# 矢印
add_arrow_right(slide, Inches(3.6), Inches(2.5), Inches(1.5), Inches(0.35), RED_PRICE)
add_arrow_right(slide, Inches(3.6), Inches(5.2), Inches(1.5), Inches(0.35), BLUE_PRICE)

# 右向き矢印（ハイハイタウン→施設）
arrow_r1 = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, Inches(7.2), Inches(2.5), Inches(1.5), Inches(0.35)
)
arrow_r1.fill.solid()
arrow_r1.fill.fore_color.rgb = ACCENT_GOLD
arrow_r1.line.fill.background()

arrow_r2 = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, Inches(7.2), Inches(5.2), Inches(1.5), Inches(0.35)
)
arrow_r2.fill.solid()
arrow_r2.fill.fore_color.rgb = ACCENT_TEAL
arrow_r2.line.fill.background()

# 下部: 所要時間
add_rounded_rect(slide, Inches(0.5), Inches(6.3), Inches(3.8), Inches(0.9),
                 LIGHT_NAVY,
                 "🚃 難波まで 近鉄 約5分", font_size=14, text_color=PRIMARY_NAVY,
                 bold=True)
add_rounded_rect(slide, Inches(4.6), Inches(6.3), Inches(3.8), Inches(0.9),
                 LIGHT_NAVY,
                 "🚃 天王寺まで 地下鉄 約5分", font_size=14, text_color=PRIMARY_NAVY,
                 bold=True)
add_rounded_rect(slide, Inches(8.7), Inches(6.3), Inches(4.1), Inches(0.9),
                 LIGHT_NAVY,
                 "🚃 梅田まで 地下鉄 約15分", font_size=14, text_color=PRIMARY_NAVY,
                 bold=True)


# ============================================================
# スライド5: 建物スペック
# ============================================================
slide = new_slide()
add_header_bar(slide, "📐 建物スペック", slide_num_counter)

# テーブル
rows = 7
cols = 2
table_shape = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.4),
                                      Inches(10), Inches(4.2))
table = table_shape.table

table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(6.5)

data = [
    ("項目", "内容"),
    ("構造", "鉄骨鉄筋コンクリート造（SRC造）"),
    ("築年月", "1980年9月（築約45年）"),
    ("総戸数", "約287戸"),
    ("土地権利", "所有権"),
    ("管理体制", "日勤管理"),
    ("間取り", "1LDK〜3LDK（約25㎡〜100㎡）"),
]

for row_idx, (col1, col2) in enumerate(data):
    for col_idx, val in enumerate([col1, col2]):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_NAME
            paragraph.font.size = Pt(14)
            if row_idx == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = DARK_TEXT
            paragraph.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT

        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_EVEN
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ODD

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# 補足
add_text_box(slide, Inches(1.5), Inches(5.9), Inches(10), Inches(0.6),
             "💡 SRC造とは？ → 鉄骨と鉄筋コンクリートを組み合わせた、とても丈夫な建て方。高層マンションによく使われます。",
             font_size=12, color=ACCENT_TEAL)


# ============================================================
# スライド6: 周辺施設マップ風
# ============================================================
slide = new_slide()
add_header_bar(slide, "🗺️ 周辺施設マップ", slide_num_counter)

# 中央建物
add_rounded_rect(slide, Inches(4.8), Inches(2.8), Inches(3.5), Inches(2.0),
                 PRIMARY_NAVY, "🏠 ハイハイタウン\n（100店舗以上の商業施設＋住居）",
                 font_size=14, text_color=WHITE, bold=True)

# 周辺施設カード
facilities = [
    ("🏬 近鉄百貨店", "隣接", Inches(0.5), Inches(1.5), ACCENT_GOLD),
    ("🏢 上本町YUFURA", "隣接", Inches(0.5), Inches(3.0), ACCENT_TEAL),
    ("🏥 大阪赤十字病院", "徒歩10分", Inches(0.5), Inches(4.5), BLUE_PRICE),
    ("🏫 学校（複数）", "徒歩圏内", Inches(0.5), Inches(6.0), RGBColor(0x7B, 0x1F, 0xA2)),
    ("🏪 コンビニ", "徒歩1分", Inches(9.5), Inches(1.5), ACCENT_TEAL),
    ("🛒 スーパー", "徒歩3分", Inches(9.5), Inches(3.0), RGBColor(0x2E, 0x7D, 0x32)),
    ("🍜 飲食店100店舗+", "館内", Inches(9.5), Inches(4.5), RED_PRICE),
    ("🚃 2路線利用可", "直結〜徒歩3分", Inches(9.5), Inches(6.0), PRIMARY_NAVY),
]

for label, dist, x, y, color in facilities:
    add_rounded_rect(slide, x, y, Inches(3.2), Inches(1.0), color,
                     f"{label}\n{dist}", font_size=12, text_color=WHITE, bold=True)


# ============================================================
# スライド7: 現在の売出し状況
# ============================================================
slide = new_slide()
add_header_bar(slide, "📊 現在の売出し状況（2026年3月）", slide_num_counter)

# フィルタ: ハイハイタウン関連（上本町駅徒歩5分以内 or 1980年9月築）
hihi_listings = [l for l in listings if
                 "上本町" in l.get("station_access", "") or
                 l.get("building_age") == "1980年9月"]

# サマリーカード
add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(1.2),
                 ACCENT_TEAL,
                 f"📦 掲載物件数\n{summary['haihaitown_count']}件",
                 font_size=16, text_color=WHITE, bold=True)

prices = [l["price_man"] for l in listings]
add_rounded_rect(slide, Inches(4.7), Inches(1.3), Inches(3.8), Inches(1.2),
                 ACCENT_GOLD,
                 f"💰 価格レンジ\n{min(prices):.0f}万〜{max(prices):.0f}万円",
                 font_size=16, text_color=WHITE, bold=True)

add_rounded_rect(slide, Inches(8.9), Inches(1.3), Inches(3.9), Inches(1.2),
                 PRIMARY_NAVY,
                 f"📈 平均㎡単価\n{summary['haihaitown_avg_price_per_sqm']:.1f}万円/㎡",
                 font_size=16, text_color=WHITE, bold=True)

# 主要物件テーブル (top 8 by relevance)
sample = sorted(listings, key=lambda x: x["price_man"])[:8]
rows = len(sample) + 1
cols = 4
tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.8),
                                    Inches(12.3), Inches(4.2))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(2.5)
tbl.columns[3].width = Inches(5.3)

headers = ["間取り", "価格（万円）", "面積（㎡）", "最寄り駅"]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER_BG
    for p in cell.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

for ri, item in enumerate(sample):
    vals = [
        item["layout"],
        f'{item["price_man"]:.0f}',
        f'{item["area_sqm"]:.1f}',
        item["station_access"][:30],
    ]
    for ci, v in enumerate(vals):
        cell = tbl.cell(ri + 1, ci)
        cell.text = v
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_ROW_EVEN if ri % 2 == 0 else TABLE_ROW_ODD
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_NAME
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# ============================================================
# スライド8: 価格帯分布
# ============================================================
slide = new_slide()
add_header_bar(slide, "📊 価格帯分布（図解）", slide_num_counter)

# 価格帯分類
ranges_def = [
    ("〜2000万", 0, 2000),
    ("2000〜3000万", 2000, 3000),
    ("3000〜4000万", 3000, 4000),
    ("4000〜5000万", 4000, 5000),
    ("5000〜7000万", 5000, 7000),
    ("7000万〜", 7000, 999999),
]

range_counts = []
for label, lo, hi in ranges_def:
    count = sum(1 for l in listings if lo <= l["price_man"] < hi)
    range_counts.append((label, count))

max_count = max(c for _, c in range_counts) if range_counts else 1

bar_colors = [BLUE_PRICE, ACCENT_TEAL, RGBColor(0x2E, 0x7D, 0x32),
              ACCENT_GOLD, RED_PRICE, RGBColor(0x7B, 0x1F, 0xA2)]

bar_start_y = Inches(1.6)
bar_height = Inches(0.65)
bar_gap = Inches(0.15)
label_width = Inches(2.5)
max_bar_width = Inches(7.5)

for i, ((label, count), color) in enumerate(zip(range_counts, bar_colors)):
    y = bar_start_y + (bar_height + bar_gap) * i

    # ラベル
    add_text_box(slide, Inches(0.5), y, label_width, bar_height,
                 label, font_size=14, color=DARK_TEXT, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    # バー
    bar_w = max(Inches(0.3), max_bar_width * (count / max_count)) if max_count > 0 else Inches(0.3)
    bar = add_rounded_rect(slide, Inches(3.2), y + Inches(0.08),
                           bar_w, bar_height - Inches(0.16),
                           color, "", font_size=12, text_color=WHITE)

    # 件数ラベル
    add_text_box(slide, Inches(3.2) + bar_w + Inches(0.2), y,
                 Inches(1.5), bar_height,
                 f"{count}件", font_size=16, color=color, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             "💡 上本町エリアは3000〜5000万円台の物件が多く、ファミリー層に人気のエリアです",
             font_size=12, color=ACCENT_TEAL)


# ============================================================
# スライド9: ㎡単価の分析
# ============================================================
slide = new_slide()
add_header_bar(slide, "💰 ㎡単価の分析", slide_num_counter)

# 「㎡単価って何？」解説ボックス
add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(1.8),
                 LIGHT_TEAL,
                 "", font_size=12)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.4),
             "💡 ㎡単価（へいべいたんか）って何？", font_size=16, bold=True,
             color=PRIMARY_NAVY)
add_multiline_text(slide, Inches(0.8), Inches(1.9), Inches(5.3), Inches(1.0),
                   ["1㎡あたりの値段のこと。",
                    "物件の大きさが違っても、㎡単価で比べると",
                    "「割高か割安か」がすぐわかります！"],
                   font_size=13, color=DARK_TEXT)

# 比較図
add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
             "📊 ㎡単価の比較", font_size=16, bold=True, color=PRIMARY_NAVY)

comparisons = [
    ("ハイハイタウン\n（上本町エリア）", f"{summary['haihaitown_avg_price_per_sqm']:.1f}万円/㎡", ACCENT_GOLD, Inches(1.9)),
    ("上本町エリア平均", "約60〜67万円/㎡", ACCENT_TEAL, Inches(3.2)),
    ("大阪市中古\nマンション平均", "約55万円/㎡", BLUE_PRICE, Inches(4.5)),
]

for label, value, color, y in comparisons:
    add_rounded_rect(slide, Inches(7), y, Inches(3.0), Inches(1.0),
                     color, label, font_size=12, text_color=WHITE, bold=True)
    add_rounded_rect(slide, Inches(10.2), y, Inches(2.5), Inches(1.0),
                     WHITE, value, font_size=16, text_color=color, bold=True)

# 解説
add_rounded_rect(slide, Inches(0.5), Inches(3.5), Inches(5.8), Inches(3.5),
                 WHITE)
add_multiline_text(slide, Inches(0.8), Inches(3.7), Inches(5.3), Inches(3.2),
                   ["📌 ポイント解説",
                    "",
                    "ハイハイタウンの㎡単価は約81.4万円/㎡で、",
                    "エリア平均（60〜67万円）より高めです。",
                    "",
                    "理由:",
                    "・近鉄大阪上本町駅に直結（駅直結プレミアム）",
                    "・商業施設併設で利便性が高い",
                    "・所有権物件（借地権ではない）",
                    "",
                    "→ 立地の良さが価格に反映されています"],
                   font_size=12, color=DARK_TEXT, line_spacing=1.3)


# ============================================================
# スライド10: 築年数と価格の関係
# ============================================================
slide = new_slide()
add_header_bar(slide, "📅 築年数と価格の関係", slide_num_counter)

# 図解: 築年数ごとの㎡単価（ビジュアルバー）
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
             "築年数が古いほど㎡単価は安くなる傾向がありますが、立地の良さが価格を支えます",
             font_size=14, color=DARK_TEXT)

age_groups = [
    ("築5年以内\n（2021年〜）", "130〜153万円/㎡", Inches(8.5), RED_PRICE),
    ("築20〜25年\n（2000年前後）", "74〜94万円/㎡", Inches(5.5), ACCENT_GOLD),
    ("築40年以上\n（1984年以前）", "42〜70万円/㎡", Inches(3.5), ACCENT_TEAL),
    ("ハイハイタウン\n（1980年・築45年）", "54〜79万円/㎡", Inches(4.0), PRIMARY_NAVY),
]

bar_x = Inches(4.5)
for i, (label, price_range, bar_width, color) in enumerate(age_groups):
    y = Inches(2.2) + Inches(1.2) * i

    add_rounded_rect(slide, Inches(0.5), y, Inches(3.5), Inches(0.9),
                     WHITE, label, font_size=12, text_color=DARK_TEXT, bold=True)

    bar = add_rounded_rect(slide, bar_x, y + Inches(0.1), bar_width, Inches(0.7),
                           color, "", font_size=11, text_color=WHITE)

    add_text_box(slide, bar_x + bar_width + Inches(0.2), y + Inches(0.15),
                 Inches(3), Inches(0.6),
                 price_range, font_size=14, color=color, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

# 注釈
add_rounded_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9),
                 LIGHT_GOLD,
                 "💡 ハイハイタウンは築45年超ですが、駅直結という立地プレミアムがあるため、同年代の物件と比べて価格が維持されています",
                 font_size=12, text_color=DARK_TEXT)


# ============================================================
# スライド11: 賃貸利回り
# ============================================================
slide = new_slide()
add_header_bar(slide, "💹 賃貸利回り", slide_num_counter)

# 利回りとは？
add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.0),
                 LIGHT_TEAL)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.4),
             "💡 利回り（りまわり）って何？", font_size=16, bold=True,
             color=PRIMARY_NAVY)
add_multiline_text(slide, Inches(0.8), Inches(1.9), Inches(5.3), Inches(1.2),
                   ["マンションを買って人に貸したとき、",
                    "1年間でもらえる家賃が、購入価格の何%か",
                    "を表す数字です。",
                    "数字が大きいほど、投資として効率が良い！"],
                   font_size=13, color=DARK_TEXT, line_spacing=1.3)

# 数値カード
add_rounded_rect(slide, Inches(7), Inches(1.3), Inches(2.8), Inches(2.0),
                 ACCENT_GOLD,
                 "想定利回り\n約4.55%", font_size=20, text_color=WHITE, bold=True)

add_rounded_rect(slide, Inches(10.1), Inches(1.3), Inches(2.7), Inches(2.0),
                 ACCENT_TEAL,
                 "平均家賃\n約8,522円/坪", font_size=18, text_color=WHITE, bold=True)

# 計算例
add_rounded_rect(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.2),
                 WHITE)
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5),
             "📝 具体例でわかる利回り計算", font_size=16, bold=True,
             color=PRIMARY_NAVY)

add_multiline_text(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.3),
                   ["【例】3,990万円で購入した2LDK（73.87㎡）を貸す場合",
                    "",
                    "  ① 坪数を計算:  73.87㎡ ÷ 3.3058 ≒ 22.35坪",
                    "  ② 月額家賃を計算:  8,522円 × 22.35坪 ≒ 190,467円/月",
                    "  ③ 年間家賃を計算:  190,467円 × 12ヶ月 ≒ 2,285,604円/年",
                    "  ④ 利回りを計算:  2,285,604円 ÷ 39,900,000円 × 100 ≒ 5.7%",
                    "",
                    "→ 毎月約19万円の家賃収入が見込めます 💰"],
                   font_size=12, color=DARK_TEXT, line_spacing=1.3)


# ============================================================
# スライド12: エリアの価格トレンド
# ============================================================
slide = new_slide()
add_header_bar(slide, "📈 エリアの価格トレンド", slide_num_counter)

# トレンドサマリー
add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(1.5),
                 RED_PRICE,
                 "上本町駅周辺\n前年比 +2.14%📈", font_size=18, text_color=WHITE, bold=True)

add_rounded_rect(slide, Inches(4.7), Inches(1.3), Inches(3.8), Inches(1.5),
                 ACCENT_GOLD,
                 "大阪市全体\nマンション価格\n上昇傾向📈", font_size=16, text_color=WHITE, bold=True)

add_rounded_rect(slide, Inches(8.9), Inches(1.3), Inches(3.9), Inches(1.5),
                 BLUE_PRICE,
                 "2025〜2026年\n市場は堅調", font_size=18, text_color=WHITE, bold=True)

# トレンド解説
add_rounded_rect(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(3.8),
                 WHITE)

trend_lines = [
    "📊 上本町エリアの市場概況",
    "",
    "● 大阪市内の中古マンション価格は2020年以降、継続的に上昇しています",
    "",
    "● 上本町エリアは、近鉄・大阪メトロの交通利便性が高く、",
    "  教育環境も良いことから、ファミリー層に根強い人気があります",
    "",
    "● 2025年の大阪万博開催を契機に、大阪市内全体の不動産需要が活発化",
    "",
    "● 上本町駅周辺の㎡単価は前年比 +2.14% で推移",
    "  → 都心部ほどの急激な上昇はなく、安定した成長を見せています",
    "",
    "● 今後も交通利便性と生活環境の良さから、堅調な推移が見込まれます",
]

add_multiline_text(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(3.5),
                   trend_lines, font_size=12, color=DARK_TEXT, line_spacing=1.2)


# ============================================================
# スライド13: メリット
# ============================================================
slide = new_slide()
add_header_bar(slide, "✅ ハイハイタウンのメリット", slide_num_counter)

merits = [
    ("🚃", "駅直結", "近鉄大阪上本町駅に直結。\n雨の日も濡れずに帰宅できます"),
    ("🏪", "商業施設併設", "1階〜4階に100店舗以上。\n日常の買い物に困りません"),
    ("📜", "所有権", "借地権ではなく所有権。\n将来の資産価値が安定"),
    ("👨‍💼", "管理体制", "日勤管理で安心。\n管理組合がしっかり運営"),
    ("🚉", "2路線利用可", "近鉄と大阪メトロの2路線。\n難波・天王寺・梅田へ快適アクセス"),
]

for i, (icon, title, desc) in enumerate(merits):
    x = Inches(0.5) + Inches(2.5) * i
    y = Inches(1.5)

    add_circle(slide, x + Inches(0.65), y, Inches(1.1), ACCENT_TEAL,
               icon, font_size=28)
    add_text_box(slide, x, y + Inches(1.3), Inches(2.3), Inches(0.5),
                 title, font_size=15, bold=True, color=PRIMARY_NAVY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.9), Inches(2.3), Inches(1.5),
                 desc, font_size=11, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)

# 下部まとめ
add_rounded_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5),
                 LIGHT_TEAL,
                 "💡 まとめ：「住む」にも「投資」にも魅力的な立地条件。駅直結＋商業併設＋所有権は、上本町エリアでも希少な組み合わせです",
                 font_size=14, text_color=PRIMARY_NAVY, bold=True)


# ============================================================
# スライド14: デメリット・注意点
# ============================================================
slide = new_slide()
add_header_bar(slide, "⚠️ デメリット・注意点", slide_num_counter)

demerits = [
    ("築45年超", "1980年竣工のため旧耐震基準。\n耐震診断・補強の有無を確認が必須"),
    ("大規模修繕", "これまでの修繕履歴と\n今後の計画・積立金残高を確認"),
    ("管理費・修繕積立金", "築年数が古い物件は\n修繕積立金が高額になる傾向"),
    ("駐車場の空き", "タワー型の場合は\n空き状況と利用料を事前確認"),
]

for i, (title, desc) in enumerate(demerits):
    y = Inches(1.5) + Inches(1.4) * i

    # 番号サークル
    add_circle(slide, Inches(0.8), y + Inches(0.1), Inches(0.6), RED_PRICE,
               str(i + 1), font_size=16)

    # カード
    add_rounded_rect(slide, Inches(1.8), y, Inches(4.5), Inches(1.15),
                     WHITE, "", font_size=12)
    add_text_box(slide, Inches(2.0), y + Inches(0.05), Inches(4.0), Inches(0.35),
                 f"⚠️ {title}", font_size=14, bold=True, color=RED_PRICE)
    add_text_box(slide, Inches(2.0), y + Inches(0.45), Inches(4.0), Inches(0.65),
                 desc, font_size=11, color=DARK_TEXT)

# 右側: 対策アドバイス
add_rounded_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.5),
                 LIGHT_NAVY)
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.0), Inches(0.5),
             "🛡️ 対策・リスク軽減のポイント", font_size=16, bold=True,
             color=PRIMARY_NAVY)

advice_lines = [
    "① 管理組合の議事録を取り寄せ、",
    "   修繕計画の進捗を確認する",
    "",
    "② 耐震診断が実施済みかどうか、",
    "   また補強工事の実績を確認する",
    "",
    "③ 修繕積立金の月額と残高を確認。",
    "   将来の値上げ予定もチェック",
    "",
    "④ リフォーム済み物件を選ぶと、",
    "   入居後の追加費用を抑えられる",
    "",
    "⑤ 複数の不動産会社に相談し、",
    "   物件の市場評価を比較する",
]
add_multiline_text(slide, Inches(7.3), Inches(2.2), Inches(5.0), Inches(4.5),
                   advice_lines, font_size=12, color=DARK_TEXT, line_spacing=1.2)


# ============================================================
# スライド15: 2030年に向けた戦略（タイムライン）
# ============================================================
slide = new_slide()
add_header_bar(slide, "🎯 2030年に向けた戦略", slide_num_counter)

# タイムライン
timeline_items = [
    ("2026", "市場監視開始\nデータ蓄積", BLUE_PRICE),
    ("2027", "狙い目の間取り\n階数を絞る", ACCENT_TEAL),
    ("2028", "資金計画\nの確定", ACCENT_GOLD),
    ("2029", "購入タイミング\nを見極め", RED_PRICE),
    ("2030", "物件取得！", PRIMARY_NAVY),
]

# 水平線（タイムライン軸）
timeline_y = Inches(3.5)
line_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.0), timeline_y + Inches(0.35),
    Inches(11.3), Inches(0.08)
)
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = MEDIUM_GRAY
line_shape.line.fill.background()

for i, (year, desc, color) in enumerate(timeline_items):
    x = Inches(1.2) + Inches(2.3) * i

    # 円（年）
    add_circle(slide, x + Inches(0.2), timeline_y - Inches(0.1),
               Inches(0.9), color, year, font_size=14, text_color=WHITE)

    # 説明ボックス（交互に上下配置）
    if i % 2 == 0:
        desc_y = timeline_y - Inches(2.0)
        # 下向き矢印的なライン
        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(0.6), timeline_y - Inches(0.9),
            Inches(0.04), Inches(0.8)
        )
    else:
        desc_y = timeline_y + Inches(1.2)
        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(0.6), timeline_y + Inches(0.45),
            Inches(0.04), Inches(0.75)
        )
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()

    add_rounded_rect(slide, x - Inches(0.2), desc_y, Inches(2.0), Inches(1.1),
                     color, desc, font_size=12, text_color=WHITE, bold=True)

# 下部コメント
add_rounded_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9),
                 LIGHT_GOLD,
                 "💡 焦らず計画的に。4年間のデータ蓄積で「最適なタイミング」が見えてきます",
                 font_size=14, text_color=DARK_TEXT, bold=True)


# ============================================================
# スライド16: 購入時のチェックポイント
# ============================================================
slide = new_slide()
add_header_bar(slide, "✅ 購入時のチェックポイント", slide_num_counter)

checks = [
    ("📋", "管理組合の議事録", "過去の議事録を確認し、トラブル履歴や将来の計画をチェック"),
    ("🔧", "大規模修繕の計画", "修繕積立金の残高と、次回の大規模修繕の時期を確認"),
    ("🏠", "リフォーム履歴", "水回り・電気設備の更新状況。リフォーム済みだと安心"),
    ("🏗️", "耐震診断・補強", "旧耐震基準の物件なので、耐震診断結果と補強の有無を確認"),
    ("🏙️", "周辺の再開発計画", "大阪市の都市計画をチェック。再開発があると資産価値UP"),
]

for i, (icon, title, desc) in enumerate(checks):
    y = Inches(1.4) + Inches(1.15) * i

    # アイコン
    add_circle(slide, Inches(0.8), y + Inches(0.05), Inches(0.8), ACCENT_TEAL,
               icon, font_size=20)

    # カード
    add_rounded_rect(slide, Inches(2.0), y, Inches(10.5), Inches(0.95),
                     WHITE)
    add_text_box(slide, Inches(2.3), y + Inches(0.05), Inches(9.8), Inches(0.35),
                 title, font_size=15, bold=True, color=PRIMARY_NAVY)
    add_text_box(slide, Inches(2.3), y + Inches(0.45), Inches(9.8), Inches(0.45),
                 desc, font_size=12, color=DARK_TEXT)


# ============================================================
# スライド17: 出典・参考情報
# ============================================================
slide = new_slide()
add_header_bar(slide, "📚 出典・参考情報", slide_num_counter)

sources = [
    ("SUUMO（スーモ）", "https://suumo.jp", "物件情報・価格データ"),
    ("マンションレビュー", "https://www.mansion-review.jp", "口コミ・評価・利回り情報"),
    ("マンションマーケット", "https://mansion-market.com", "㎡単価・価格トレンド"),
    ("LIFULL HOME'S 住まいインデックス",
     "https://lifullhomes-index.jp", "エリア価格推移データ"),
    ("国土交通省 不動産情報ライブラリ",
     "https://www.reinfolib.mlit.go.jp", "公示地価・取引価格"),
]

for i, (name, url, desc) in enumerate(sources):
    y = Inches(1.5) + Inches(1.05) * i

    add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.85),
                     WHITE)

    add_circle(slide, Inches(1.0), y + Inches(0.1), Inches(0.6), ACCENT_TEAL,
               str(i + 1), font_size=14)

    add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(4.5), Inches(0.35),
                 name, font_size=14, bold=True, color=PRIMARY_NAVY)
    add_text_box(slide, Inches(6.5), y + Inches(0.05), Inches(5.5), Inches(0.35),
                 url, font_size=12, color=BLUE_PRICE)
    add_text_box(slide, Inches(1.8), y + Inches(0.42), Inches(10), Inches(0.35),
                 desc, font_size=11, color=MEDIUM_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
             "※ 掲載データは2026年3月時点の情報です。最新情報は各サイトでご確認ください。",
             font_size=11, color=MEDIUM_GRAY)


# ============================================================
# スライド18: まとめ
# ============================================================
slide = new_slide()
add_header_bar(slide, "📝 まとめ ― 美玖のコメント", slide_num_counter)

# メインコメントカード
add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.8),
                 WHITE)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
             "💬 美玖の総合評価", font_size=18, bold=True, color=PRIMARY_NAVY)

comment_lines = [
    "上本町ハイハイタウンは、築45年超の物件ですが、「駅直結」「商業施設併設」「所有権」",
    "という3つの大きな強みを持つ、上本町エリアでも希少な物件です。",
    "",
    "㎡単価はエリア平均よりやや高めですが、それは立地プレミアムの証。",
    "特に近鉄大阪上本町駅に直結している利便性は、将来の資産価値を支える大きな要因です。",
    "",
    "2030年の購入に向けて、今から市場を監視し、データを蓄積していくことをおすすめします。",
]
add_multiline_text(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.0),
                   comment_lines, font_size=13, color=DARK_TEXT, line_spacing=1.3)

# キーポイント3つ
key_points = [
    ("立地", "駅直結+2路線\n利用可能", ACCENT_TEAL),
    ("価格", "㎡単価 81.4万円\nエリア平均より高め\nだが立地で納得", ACCENT_GOLD),
    ("戦略", "2030年目標で\n4年間じっくり\n準備する", PRIMARY_NAVY),
]

for i, (title, desc, color) in enumerate(key_points):
    x = Inches(0.5) + Inches(4.2) * i
    add_rounded_rect(slide, x, Inches(4.5), Inches(3.8), Inches(2.5),
                     color,
                     f"🔑 {title}\n\n{desc}",
                     font_size=14, text_color=WHITE, bold=True)


# ============================================================
# 保存
# ============================================================
output_path = os.path.join(SCRIPT_DIR, "上本町ハイハイタウン_エリア分析レポート_202603.pptx")
prs.save(output_path)
print(f"パワーポイントを生成しました！")
print(f"ファイルパス: {output_path}")
print(f"スライド数: {slide_num_counter}枚")
